from __future__ import annotations

import base64
import io
import json
import logging
from dataclasses import dataclass
from typing import Any

import httpx
try:
    from docx import Document as DocxDocument
except ImportError:  # pragma: no cover - optional until environment is provisioned
    DocxDocument = None
try:
    from pptx import Presentation  # type: ignore[import-not-found]
except ImportError:  # pragma: no cover - optional until environment is provisioned
    Presentation = None
try:
    import pytesseract  # type: ignore[import-not-found]
except ImportError:  # pragma: no cover - optional until environment is provisioned
    pytesseract = None
from PIL import Image
from pypdf import PdfReader

from app.core.config import settings

logger = logging.getLogger("quizify.extraction")

if settings.tesseract_cmd and pytesseract is not None:
    pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd


@dataclass
class ExtractionResult:
    extracted_text: str
    chunks: list[str]
    diagram_caption: str | None
    metadata: dict[str, Any]


class ExtractionService:
    def extract(self, file_name: str, source_type: str, payload: bytes, fallback_text: str = "") -> ExtractionResult:
        normalized = (source_type or "").lower()
        if normalized == "auto":
            normalized = self._detect_source_type(file_name)

        text = ""
        caption = None
        if normalized == "pdf":
            text = self._extract_pdf_text(payload)
        elif normalized == "docx":
            text = self._extract_docx_text(payload)
        elif normalized == "pptx":
            text = self._extract_pptx_text(payload)
        elif normalized in {"image", "camera"}:
            text = self._extract_image_text(payload)
            caption = self._caption_diagram(payload)
            if caption:
                text = f"{text}\n\nDiagram interpretation: {caption}".strip()
        elif normalized in {"text", "txt", "markdown", "csv"}:
            text = self._extract_plain_text(payload)
        else:
            text = fallback_text.strip()

        fallback = fallback_text.strip()
        text = self._normalize_text(text)
        if len(text) < 120 and fallback:
            # If extraction is too thin, append user-provided fallback for better quiz quality.
            text = f"{text}\n\n{fallback}".strip()
        elif not text:
            text = fallback

        if not text.strip():
            raise ValueError(
                "No extractable text found. Upload a PDF/image with readable text or add fallback text."
            )

        chunks = self._chunk_text(text)
        return ExtractionResult(
            extracted_text=text,
            chunks=chunks,
            diagram_caption=caption,
            metadata={
                "source_type": normalized,
                "file_name": file_name,
                "chunk_count": len(chunks),
                "has_diagram_caption": bool(caption),
            },
        )

    def _detect_source_type(self, file_name: str) -> str:
        lowered = file_name.lower()
        if lowered.endswith(".pdf"):
            return "pdf"
        if lowered.endswith(".docx"):
            return "docx"
        if lowered.endswith(".pptx"):
            return "pptx"
        if lowered.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff")):
            return "image"
        if lowered.endswith((".txt", ".md", ".csv")):
            return "text"
        return "text"

    def _extract_plain_text(self, payload: bytes) -> str:
        try:
            return payload.decode("utf-8").strip()
        except UnicodeDecodeError:
            try:
                return payload.decode("latin-1").strip()
            except UnicodeDecodeError:
                logger.warning("Plain text extraction failed: unsupported encoding")
                return ""

    def _extract_docx_text(self, payload: bytes) -> str:
        if DocxDocument is None:
            logger.warning("DOCX extraction skipped: python-docx not installed")
            return ""

        try:
            doc = DocxDocument(io.BytesIO(payload))
            lines = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip()]
            return "\n".join(lines).strip()
        except Exception as exc:
            logger.warning("DOCX extraction failed: %s", exc)
            return ""

    def _extract_pdf_text(self, payload: bytes) -> str:
        try:
            reader = PdfReader(io.BytesIO(payload))
            output: list[str] = []
            for page in reader.pages:
                text = page.extract_text() or ""
                if text.strip():
                    output.append(text.strip())
            return "\n\n".join(output)
        except Exception as exc:
            logger.warning("PDF extraction failed: %s", exc)
            return ""

    def _extract_pptx_text(self, payload: bytes) -> str:
        if Presentation is None:
            logger.warning("PPTX extraction skipped: python-pptx not installed")
            return ""

        try:
            deck = Presentation(io.BytesIO(payload))
            lines: list[str] = []
            for slide_index, slide in enumerate(deck.slides, start=1):
                slide_header_added = False
                for shape in slide.shapes:
                    raw = getattr(shape, "text", "") or ""
                    normalized = raw.strip()
                    if not normalized:
                        continue
                    if not slide_header_added:
                        lines.append(f"Slide {slide_index}:")
                        slide_header_added = True
                    lines.append(normalized)
            return "\n".join(lines).strip()
        except Exception as exc:
            logger.warning("PPTX extraction failed: %s", exc)
            return ""

    def _extract_image_text(self, payload: bytes) -> str:
        if pytesseract is None:
            return ""
        try:
            image = Image.open(io.BytesIO(payload))
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as exc:
            logger.warning("Image extraction failed: %s", exc)
            return ""

    def _caption_diagram(self, payload: bytes) -> str | None:
        if not settings.azure_openai_api_key or not settings.azure_openai_endpoint:
            return None

        image_b64 = base64.b64encode(payload).decode("utf-8")
        endpoint = (
            f"{settings.azure_openai_endpoint}openai/deployments/"
            f"{settings.azure_openai_deployment_chat}/chat/completions"
            f"?api-version={settings.azure_openai_api_version}"
        )
        headers = {
            "api-key": settings.azure_openai_api_key,
            "Content-Type": "application/json",
        }
        body = {
            "messages": [
                {
                    "role": "system",
                    "content": "You are a visual educational assistant. Describe diagrams, charts, and flow direction with brief educational caption.",
                },
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Caption this educational image/diagram in 4-6 lines."},
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/png;base64,{image_b64}"},
                        },
                    ],
                },
            ],
            "temperature": 0.2,
        }

        try:
            with httpx.Client(timeout=45.0) as client:
                response = client.post(endpoint, headers=headers, json=body)
                response.raise_for_status()
                data = response.json()
            return data["choices"][0]["message"]["content"].strip()
        except (httpx.HTTPError, KeyError, IndexError, json.JSONDecodeError):
            return None

    def _chunk_text(self, text: str, chunk_size: int = 900, overlap: int = 150) -> list[str]:
        if not text.strip():
            return []

        chunks: list[str] = []
        cursor = 0
        total = len(text)
        while cursor < total:
            end = min(total, cursor + chunk_size)
            chunk = text[cursor:end].strip()
            if chunk:
                chunks.append(chunk)
            if end == total:
                break
            cursor = max(0, end - overlap)
        return chunks

    def _normalize_text(self, text: str) -> str:
        if not text:
            return ""
        cleaned = text.replace("\x00", " ").replace("\r\n", "\n").replace("\r", "\n")
        cleaned = "\n".join(line.strip() for line in cleaned.split("\n"))
        cleaned = "\n".join(line for line in cleaned.split("\n") if line)
        return cleaned.strip()
